"""
Studio artifact generation service.
Generates 9 types of AI content from a project's documents:
  podcast, slides, video_script, mindmap, report,
  flashcards, quiz, infographic, datatable
"""
import asyncio
import json
import logging
import shutil
import tempfile
from pathlib import Path

from llama_index.core.llms import ChatMessage, MessageRole

from app.services.llm_service import get_llm, _fresh_async_client
from app.services.summary_service import _collect_document_text
from app.models import (
    list_project_documents,
    update_studio_artifact,
    STUDIO_ARTIFACT_TYPES,
)
from pydantic import ValidationError
from app.schemas.slides import SlidesSpec

# ── System prompts ────────────────────────────────────────────

PODCAST_PROMPT = """你是一位專業的 Podcast 腳本撰寫人。請根據以下文件內容，撰寫一段雙人主持的對話腳本。

請嚴格按照以下 JSON 格式輸出（不要加入其他文字或 markdown 標記）：
{
  "host_a": ["主持人A第一句台詞", "主持人A第二句台詞", ...],
  "host_b": ["主持人B第一句台詞", "主持人B第二句台詞", ...]
}

規則：
1. 使用繁體中文。
2. 兩位主持人輪流對話，每人 8-12 句。
3. 對話自然生動，像真實的討論節目。
4. 涵蓋文件的主要重點與有趣細節。
5. 輸出必須是合法的 JSON。
"""

_SLIDES_OUTLINE_PROMPT = """\
你是簡報架構師。根據以下文件內容，先規劃一份簡報大綱，再由另一個步驟填充完整內容。

請分析文件主題，輸出 JSON 格式的大綱：
{
  "narrative": "匯報|提案|分析|教學",
  "theme": "tech|ocean|golden|frost|garden|sports",
  "chapters": [
    {
      "chapter_title": "章節標題",
      "key_point": "這章節要傳達的核心論點（一句話）",
      "suggested_layout": "建議版面類型（如 big_number、card_grid 等）",
      "data_available": true或false（文件中是否有具體數字/統計）
    }
  ]
}

選擇敘事弧線的依據：
- 匯報：文件描述現況、成效、數據 → 背景→發現→數據→啟示→結論
- 提案：文件提出解決方案 → 痛點→方案→佐證→行動
- 分析：文件分析法規、政策、概念 → 定義→拆解→利弊→建議
- 教學：文件解釋知識、技術、流程 → 重要性→概念→步驟→回顧

選擇主題色票的依據：
- tech：科技、AI、軟體、數位轉型
- ocean：環境、醫療、教育、公共政策、法律
- golden：金融、商業、行銷、品牌
- frost：學術、研究、法律白皮書（嚴謹正式）
- garden：農業、食品、永續、ESG
- sports：體育、賽事、健康、活力

規則：
- chapters 數量 12-15 個（對應 12-15 頁內容，不含封面和結尾）
- 嚴格按照敘事弧線的邏輯順序排列 chapters
- suggested_layout 從以下選擇：big_number、card_grid、dual_column、process_flow、content_with_icon、quote_slide、table、chart、section_divider
- 只輸出 JSON，不加說明文字
"""

SLIDES_PROMPT = """你是專業簡報設計師。根據文件內容輸出 JSON 格式的簡報結構。

## 色票選擇（根據文件主題選一）
| 名稱 | 適用場景 |
|------|---------|
| tech | 科技、AI、軟體、數位轉型 |
| ocean | 環境、醫療、教育、公共政策 |
| golden | 金融、商業、行銷、品牌策略 |
| frost | 學術、研究、法律、白皮書 |
| garden | 農業、食品、永續、ESG |
| sports | 體育、賽事、健康、活力主題 |

## 可用版面類型
- cover: 封面（title, subtitle）
- section_divider: 章節分隔（label, title, description）
- big_number: 關鍵指標（title, items:[{value,unit,label}], 1-3個）
- card_grid: 並列卡片（title, cards:[{icon,title,description}], 2-4張）
- dual_column: 比較（title, left/right:{icon,title,points:[]}）
- process_flow: 流程步驟（title, steps:[{title,description}], 2-5步）
- content_with_icon: 圖文（title, icon, blocks:[{title,description}], 1-4個）⚠️ blocks 為必填，不可省略
- quote_slide: 金句（quote, source）
- table: 表格（title, headers:[], rows:[[]]）⚠️ rows 中每個子陣列長度必須等於 headers 長度，否則驗證失敗
- chart: 圖表（title, chart_type:BAR|PIE, labels:[], values:[]）⚠️ labels 與 values 長度必須相同；僅當文件含具體數字時使用
- conclusion: 結尾（title, summary, points:[{text,icon}]）

可用 icon 及語義對照（先判斷概念語義，再對應 icon）：
成長/趨勢/成果 → FaChartLine | 保護/防禦/安全 → FaShieldAlt
創意/策略/靈感 → FaLightbulb | 法律/規範/裁定 → FaGavel
資料/儲存/結構 → FaDatabase  | 全球/國際/網路 → FaGlobe
人群/團隊/組織 → FaUsers     | 驗證/完成/通過 → FaCheck
隱私/權限/加密 → FaLock      | 學習/知識/文獻 → FaBook
搜尋/分析/調查 → FaSearch    | 目標/里程碑/旗幟 → FaFlag
發射/行動/創新 → FaRocket    | 合作/協議/握手 → FaHandshake
系統/設定/機制 → FaCog        | 報表/統計/數據 → FaChartBar

Icon 選用三規則：
1. 先將概念抽象成語義類別，再從對照表選最接近的 icon
2. 同一頁多個 icon 必須在形狀上明顯不同，讓讀者掃一眼即可分辨各塊的差異
3. 每個 icon 在整份簡報最多出現 2 次

## 內容規則
- 首頁 cover，末頁 conclusion
- 相鄰兩頁不可相同 layout，全簡報 ≥5 種不同 layout
- 同一 layout 類型在全簡報最多出現 2 次；每個 layout 類型出現後，至少間隔 3 頁才可再次使用
- 繁體中文，不編造數據或引言
- 標題≤15字，要點≤25字
- 有明確數字優先用 big_number；chart 僅用於文件有具體數字時
- 敘事弧線：匯報(背景→發現→數據→啟示→結論) | 提案(痛點→方案→佐證→行動) | 分析(定義→拆解→利弊→建議) | 教學(重要性→概念→步驟→回顧)

## 文字密度規則（超出上限驗證會失敗）
投影片只放關鍵字短語；完整解釋、背景脈絡放 speaker_notes。

| 欄位 | 上限 |
|------|------|
| 所有版面 title | 20字 |
| cover subtitle | 30字 |
| section_divider description | 40字 |
| process_flow step description | 40字 |
| content_with_icon block description | 40字 |
| big_number value | 10字 |
| big_number unit / label | 10 / 15字 |
| quote_slide quote | 60字 |
| table header cell | 20字 |
| table data cell | 20字 |
| speaker_notes | 50–200字 |

description 欄位一律使用純繁體中文短語，禁止夾雜英文術語（如 stakeholder、organisational）；英文專有名詞移到 speaker_notes。

⚠️ description 欄位必須補充 title 以外的資訊，不得重複或意譯 title 內容。錯誤示範：title="坦承接受懲罰" description="坦承並接受懲罰" — 完全重複，不可接受。

## 輸出範例（4頁示意，實際需封面＋12-15頁內容＋結尾，共 14-17 頁）
```json
{
  "theme": "tech",
  "narrative": "匯報",
  "slides": [
    {"layout": "cover", "title": "AI 導入成效報告", "subtitle": "2026 Q1 季度回顧",
     "speaker_notes": "本報告摘要 2026 第一季 AI 系統導入後的主要成效，涵蓋模型準確率、處理速度提升與下一步擴展計畫。"},
    {"layout": "big_number", "title": "關鍵指標", "items": [
      {"value": "98%", "unit": "準確率", "label": "模型推論"},
      {"value": "3.2x", "unit": "加速", "label": "處理速度"}
    ], "speaker_notes": "兩項核心指標均超越原訂目標。98% 準確率是在正式環境連續兩周測試後取得的穩定值；3.2 倍加速效果已通過壓力測試驗證。"},
    {"layout": "card_grid", "title": "三大策略方向", "cards": [
      {"icon": "FaRocket", "title": "擴展部署", "description": "推廣至五個部門"},
      {"icon": "FaDatabase", "title": "資料整合", "description": "統一資料湖架構"},
      {"icon": "FaUsers", "title": "人才培訓", "description": "培訓兩百名工程師"}
    ], "speaker_notes": "三項策略並行推進：部署擴展聚焦高流量業務部門；資料整合解決現有資料孤島問題；人才培訓確保技術落地後有足夠人力維運。"},
    {"layout": "conclusion", "title": "總結與展望", "summary": "AI 導入已見初步成效",
     "points": [{"text": "模型準確率達 98%"}, {"text": "處理速度提升 3.2 倍"}, {"text": "下季擴展至全公司", "icon": "FaRocket"}],
     "speaker_notes": "本季成效驗證了 AI 導入策略的可行性。下一季將全面推廣，預期帶來更大規模的效率提升。歡迎提問。"}
  ]
}
```

## 輸出規則
- 只輸出 JSON，不加 ```json 標記或任何說明文字
- slides 陣列必須包含 14-17 頁（1 封面 ＋ 12-15 內容頁 ＋ 1 結尾）
- 充分利用文件素材，每個主題/章節至少一頁，不要將多個主題壓縮在一頁
- 每頁必須包含 speaker_notes（50–200字的繁體中文完整句子，供演講者現場參考）
- 輸出必須是合法的 JSON
"""

VIDEO_SCRIPT_PROMPT = """你是一位專業的影片旁白撰寫人。請根據以下文件內容，撰寫一段影片解說旁白腳本。

請嚴格按照以下 JSON 格式輸出（不要加入其他文字或 markdown 標記）：
{
  "script": "完整的旁白腳本文字，以段落分隔，包含場景提示用[...]標示"
}

規則：
1. 使用繁體中文。
2. 旁白長度約 400-600 字。
3. 語氣清晰、適合朗讀。
4. 可加入簡短的場景提示，例如 [顯示圖表]、[切換畫面]。
5. 涵蓋文件核心內容。
6. 輸出必須是合法的 JSON。
"""

MINDMAP_PROMPT = """你是一位知識整理專家。請根據以下文件內容，產生一份心智圖結構。

請嚴格按照以下 JSON 格式輸出（不要加入其他文字或 markdown 標記）：
{
  "root": {
    "label": "中心主題",
    "children": [
      {
        "label": "主要分支一",
        "children": [
          {"label": "子節點一"},
          {"label": "子節點二"}
        ]
      }
    ]
  }
}

規則：
1. 使用繁體中文。
2. 根節點為文件主題。
3. 第一層 4-6 個主要分支。
4. 每個分支 2-4 個子節點。
5. 第三層節點不再展開（最多三層）。
6. 輸出必須是合法的 JSON。
"""

REPORT_PROMPT = """你是一位專業的分析報告撰寫人。請根據以下文件內容，撰寫一份完整的分析報告。

請嚴格按照以下 JSON 格式輸出（不要加入其他文字或 markdown 標記）：
{
  "markdown": "完整的 Markdown 格式報告（使用 # ## ### 標題層級，- 列表，**粗體** 強調）"
}

規則：
1. 使用繁體中文。
2. 報告包含：執行摘要、背景、主要發現、分析、結論與建議。
3. 長度 600-1000 字。
4. 使用標準 Markdown 語法。
5. 輸出必須是合法的 JSON。
"""

FLASHCARDS_PROMPT = """你是一位教育內容設計師。請根據以下文件內容，製作一組學習卡片。

請嚴格按照以下 JSON 格式輸出（不要加入其他文字或 markdown 標記）：
{
  "cards": [
    {"front": "問題或概念", "back": "答案或解釋"},
    ...
  ]
}

規則：
1. 使用繁體中文。
2. 製作 10-15 張卡片。
3. 正面是問題、概念或術語。
4. 背面是簡潔的答案或定義（1-3 句）。
5. 涵蓋文件的關鍵知識點。
6. 輸出必須是合法的 JSON。
"""

QUIZ_PROMPT = """你是一位出題老師。請根據以下文件內容，設計一份選擇題測驗。

請嚴格按照以下 JSON 格式輸出（不要加入其他文字或 markdown 標記）：
{
  "questions": [
    {
      "q": "題目文字？",
      "options": ["選項A", "選項B", "選項C", "選項D"],
      "answer": 0,
      "explanation": "解析說明"
    }
  ]
}

規則：
1. 使用繁體中文。
2. 出 8-10 道選擇題。
3. 每題四個選項，answer 為正確答案的索引（0-3）。
4. 難易適中，考察理解而非死記。
5. 每題附上解析。
6. 輸出必須是合法的 JSON。
"""

INFOGRAPHIC_PROMPT = """你是一位資料視覺化專家。請根據以下文件內容，萃取出適合以圖表呈現的數據或比較資訊。

請嚴格按照以下 JSON 格式輸出（不要加入其他文字或 markdown 標記）：
{
  "chart_type": "bar",
  "title": "圖表標題",
  "labels": ["項目A", "項目B", "項目C"],
  "datasets": [
    {
      "label": "數據系列名稱",
      "data": [100, 200, 150]
    }
  ]
}

規則：
1. 使用繁體中文。
2. chart_type 從 "bar"（長條圖）、"pie"（圓餅圖）、"line"（折線圖）中選最適合的。
3. labels 3-8 個項目。
4. 若文件無明確數字，可用相對重要性（1-10 分）或出現頻率估計。
5. 輸出必須是合法的 JSON。
"""

DATATABLE_PROMPT = """你是一位資料分析師。請根據以下文件內容，萃取或整理成一份結構化的資料表格。

請嚴格按照以下 JSON 格式輸出（不要加入其他文字或 markdown 標記）：
{
  "title": "表格標題",
  "headers": ["欄位一", "欄位二", "欄位三"],
  "rows": [
    ["資料A1", "資料A2", "資料A3"],
    ["資料B1", "資料B2", "資料B3"]
  ]
}

規則：
1. 使用繁體中文。
2. 3-6 個欄位，5-15 筆資料。
3. 選擇文件中最有價值的比較或清單資訊製表。
4. 若文件含多個主題，可整理成重要概念對照表。
5. 輸出必須是合法的 JSON。
"""

ARTIFACT_PROMPTS: dict[str, str] = {
    "podcast": PODCAST_PROMPT,
    "slides": SLIDES_PROMPT,
    "video_script": VIDEO_SCRIPT_PROMPT,
    "mindmap": MINDMAP_PROMPT,
    "report": REPORT_PROMPT,
    "flashcards": FLASHCARDS_PROMPT,
    "quiz": QUIZ_PROMPT,
    "infographic": INFOGRAPHIC_PROMPT,
    "datatable": DATATABLE_PROMPT,
}

# Types where content_text holds the plain-text (copy-friendly) version
_TEXT_ONLY_TYPES = {"video_script", "report"}

_MAX_PER_DOC_CHARS = 8000
_MAX_TOTAL_CHARS = 20000
_SLIDES_MAX_PER_DOC_CHARS = 6000   # enough material per doc for 8-12 slides
_SLIDES_MAX_TOTAL_CHARS = 15000    # allow richer context so LLM can fill 8-12 pages
_PROGRESS_UPDATE_EVERY = 500  # chars between DB progress updates
_STREAM_TIMEOUT_SECS = 300    # max seconds to wait for LLM streaming (including first-token latency)
_STREAM_MAX_RETRIES = 2       # retry attempts when streaming returns 0 chars (server busy)


def _strip_code_fence(raw: str) -> str:
    """Remove markdown ```json ... ``` wrappers from LLM output."""
    if raw.startswith("```"):
        lines = [l for l in raw.split("\n") if not l.strip().startswith("```")]
        return "\n".join(lines).strip()
    return raw


def _sanitize_json(raw: str) -> str:
    """
    Clean common LLM JSON output issues before Pydantic validation.
    1. Strip markdown code fences (```json ... ```)
    2. Remove trailing commas before } or ]
    3. Unescape literal newlines inside JSON strings
    """
    import re
    # 1. Strip code fences
    text = _strip_code_fence(raw)
    # 2. Remove trailing commas: {"a": 1,} → {"a": 1}
    text = re.sub(r",\s*([}\]])", r"\1", text)
    # 3. Replace literal (unescaped) newlines inside JSON strings with \\n
    # Only replace newlines that are inside double-quoted strings
    def _fix_newlines(m: re.Match) -> str:
        return m.group(0).replace("\n", "\\n").replace("\r", "")
    text = re.sub(r'"[^"\\]*(?:\\.[^"\\]*)*"', _fix_newlines, text, flags=re.DOTALL)
    return text.strip()


async def _fix_slides_json(raw_json: str, validation_error: str) -> str:
    """
    Ask the LLM to fix a JSON validation error in slides spec.
    Returns the corrected JSON string, or the original if the fix fails.
    Uses slides_model if configured.
    """
    from app.routers.settings import _runtime_settings as _rs
    fix_prompt = (
        "以下 JSON 簡報資料有驗證錯誤，請修正並只輸出修正後的完整 JSON，不加任何說明或 markdown 標記。\n\n"
        "欄位字元上限（請一次檢查所有欄位，不只修錯誤訊息中提到的那一個）：\n"
        "- 所有版面 title：20字\n"
        "- cover subtitle：30字\n"
        "- section_divider description：40字\n"
        "- process_flow step description：40字\n"
        "- content_with_icon block description：40字\n"
        "- big_number value：10字，unit/label：10/15字\n"
        "- quote_slide quote：60字\n"
        "- table header cell：20字，data cell：20字\n"
        "- conclusion point text：25字\n"
        "- dual_column points 每條：25字\n\n"
        f"錯誤訊息：\n{validation_error}\n\n"
        f"原始 JSON：\n{raw_json}"
    )
    messages = [ChatMessage(role=MessageRole.USER, content=fix_prompt)]
    fix_client = _fresh_async_client()
    llm = get_llm(
        async_client=fix_client,
        model_override=_rs.slides_model if _rs.slides_model else None,
    )
    try:
        parts: list[str] = []
        async with asyncio.timeout(90):
            async for chunk in await llm.astream_chat(messages):
                if chunk.delta:
                    parts.append(chunk.delta)
        fixed = "".join(parts).strip()
        return _sanitize_json(fixed) if fixed else raw_json
    except asyncio.TimeoutError:
        logging.warning("_fix_slides_json timed out after 90s — using original JSON")
        return raw_json
    except Exception:
        logging.exception("_fix_slides_json LLM call failed — using original JSON")
        return raw_json
    finally:
        await fix_client.aclose()


def _format_text(artifact_type: str, data: dict) -> str:
    """Extract a plain-text representation from parsed JSON data."""
    if artifact_type == "video_script":
        return data.get("script", "")
    if artifact_type == "report":
        return data.get("markdown", "")
    if artifact_type == "podcast":
        # Interleave host_a and host_b lines
        host_a = data.get("host_a", [])
        host_b = data.get("host_b", [])
        lines = []
        for i, (a, b) in enumerate(zip(host_a, host_b)):
            lines.append(f"主持人 A：{a}")
            lines.append(f"主持人 B：{b}")
        # Append remaining lines if one list is longer
        for line in host_a[len(host_b):]:
            lines.append(f"主持人 A：{line}")
        for line in host_b[len(host_a):]:
            lines.append(f"主持人 B：{line}")
        return "\n\n".join(lines)
    return ""


# ── Core generation function ──────────────────────────────────

async def _generate_slides_outline(
    combined_text: str,
    client,
) -> dict | None:
    """
    Stage 1 of two-stage slides generation: produce a narrative outline.
    Returns a dict with keys: narrative, theme, chapters — or None on failure.
    """
    from app.services.llm_service import get_llm
    llm = get_llm(async_client=client)
    user_msg = f"請分析以下文件內容並規劃簡報大綱：\n\n{combined_text}"
    messages = [
        ChatMessage(role=MessageRole.SYSTEM, content=_SLIDES_OUTLINE_PROMPT),
        ChatMessage(role=MessageRole.USER, content=user_msg),
    ]
    try:
        parts: list[str] = []
        async with asyncio.timeout(120):
            async for chunk in await llm.astream_chat(messages):
                if chunk.delta:
                    parts.append(chunk.delta)
        raw = "".join(parts).strip()
        raw = _strip_code_fence(raw)
        outline = json.loads(raw)
        if isinstance(outline, dict) and "chapters" in outline:
            return outline
        return None
    except Exception:
        logging.warning("Slides outline generation failed — proceeding without outline")
        return None


async def generate_artifact(project_id: int, artifact_id: int, artifact_type: str) -> None:
    """
    Background task: collect all document text for the project,
    call the LLM with the appropriate prompt, then persist the result.
    """
    raw = ""
    try:
        update_studio_artifact(artifact_id, status="generating", progress_message="正在讀取文件內容…")

        # Collect text from all ready documents in the project
        docs = list_project_documents(project_id)
        parts: list[str] = []
        for doc in docs:
            if doc.status == "ready":
                text = _collect_document_text(doc.collection_name, max_chars=_MAX_PER_DOC_CHARS)
                if text:
                    parts.append(f"=== 文件：{doc.filename} ===\n{text}")

        if not parts:
            update_studio_artifact(
                artifact_id,
                status="error",
                error_message="專案中尚無可用的文件內容，請先上傳並等待文件處理完成。",
            )
            return

        if artifact_type not in ARTIFACT_PROMPTS:
            update_studio_artifact(
                artifact_id,
                status="error",
                error_message=f"不支援的 artifact 類型：{artifact_type}",
            )
            return

        # Slides use a tighter text budget to keep the total LLM context smaller,
        # which reduces model reasoning time (TTFT) significantly.
        outline = None
        if artifact_type == "slides":
            parts_slides: list[str] = []
            for doc in docs:
                if doc.status == "ready":
                    text = _collect_document_text(doc.collection_name, max_chars=_SLIDES_MAX_PER_DOC_CHARS)
                    if text:
                        parts_slides.append(f"=== 文件：{doc.filename} ===\n{text}")
            combined = "\n\n".join(parts_slides) if parts_slides else "\n\n".join(parts)
            if len(combined) > _SLIDES_MAX_TOTAL_CHARS:
                combined = combined[:_SLIDES_MAX_TOTAL_CHARS] + "\n\n…（內容已截斷）"
        else:
            combined = "\n\n".join(parts)
            if len(combined) > _MAX_TOTAL_CHARS:
                combined = combined[:_MAX_TOTAL_CHARS] + "\n\n…（內容已截斷）"

        prompt = ARTIFACT_PROMPTS[artifact_type]

        # Use a fresh AsyncClient per generation to avoid stale connection
        # pool state after asyncio cancellation (previous timeout can corrupt
        # the shared pool, causing the next streaming request to receive 0 bytes).
        _stream_client = _fresh_async_client()

        if artifact_type == "slides":
            # Two-stage generation: first generate an outline to anchor the narrative structure
            update_studio_artifact(artifact_id, progress_message="AI 正在分析文件架構，規劃敘事大綱…")
            outline = await _generate_slides_outline(combined, _stream_client)

        if artifact_type == "slides" and outline:
            outline_json = json.dumps(outline, ensure_ascii=False, indent=2)
            user_msg = (
                f"【敘事大綱（請嚴格按照此架構生成簡報）】\n{outline_json}\n\n"
                f"【文件內容】\n{combined}"
            )
        else:
            user_msg = f"以下是專案的所有文件內容：\n\n{combined}"

        # Slides may use a dedicated larger model (slides_model) to improve
        # generation quality and reduce TTFT on complex prompts.
        from app.routers.settings import _runtime_settings as _rs
        llm = get_llm(
            async_client=_stream_client,
            model_override=_rs.slides_model if artifact_type == "slides" and _rs.slides_model else None,
            max_tokens_override=8192 if artifact_type == "slides" else None,
        )
        messages = [
            ChatMessage(role=MessageRole.SYSTEM, content=prompt),
            ChatMessage(role=MessageRole.USER, content=user_msg),
        ]

        update_studio_artifact(artifact_id, progress_message="AI 正在生成內容，請耐心等候…")

        # Streaming with auto-retry on zero-char timeout.
        # The model server can be temporarily overloaded, causing 0 chars within
        # the timeout. A fresh client + brief pause before retry often succeeds.
        raw_parts: list[str] = []
        total_chars = 0
        last_progress_chars = 0
        for _attempt in range(1 + _STREAM_MAX_RETRIES):
            if _attempt > 0:
                logging.info(
                    "Retrying LLM stream for artifact %d (attempt %d/%d)",
                    artifact_id, _attempt + 1, 1 + _STREAM_MAX_RETRIES,
                )
                update_studio_artifact(
                    artifact_id,
                    progress_message=f"AI 伺服器忙碌，自動重試中（第 {_attempt + 1} 次）…",
                )
                await asyncio.sleep(5)
                await _stream_client.aclose()
                _stream_client = _fresh_async_client()
                llm = get_llm(async_client=_stream_client)

            raw_parts = []
            total_chars = 0
            last_progress_chars = 0
            try:
                async with asyncio.timeout(_STREAM_TIMEOUT_SECS):
                    async for chunk in await llm.astream_chat(messages):
                        if chunk.delta:
                            raw_parts.append(chunk.delta)
                            total_chars += len(chunk.delta)
                            if total_chars - last_progress_chars >= _PROGRESS_UPDATE_EVERY:
                                last_progress_chars = total_chars
                                update_studio_artifact(
                                    artifact_id,
                                    progress_message=f"AI 正在生成…已產生約 {total_chars} 字",
                                )
            except TimeoutError:
                logging.warning(
                    "LLM streaming timed out after %ds for artifact %d (got %d chars, attempt %d)",
                    _STREAM_TIMEOUT_SECS, artifact_id, total_chars, _attempt + 1,
                )
                if total_chars == 0 and _attempt < _STREAM_MAX_RETRIES:
                    continue  # retry
                if not raw_parts:
                    update_studio_artifact(
                        artifact_id,
                        status="error",
                        error_message=f"AI 伺服器忙碌，{_STREAM_TIMEOUT_SECS} 秒內未回應（已重試 {_attempt + 1} 次），請稍後重試。",
                    )
                    await _stream_client.aclose()
                    return
                # Got partial content — try to use what we have
            break  # success or partial — stop retrying
        await _stream_client.aclose()
        raw = "".join(raw_parts).strip()
        raw = _strip_code_fence(raw)

        if not raw:
            update_studio_artifact(
                artifact_id,
                status="error",
                error_message="AI 未產生任何內容，可能模型暫時無法回應，請稍後重試。",
            )
            return

        if artifact_type == "slides":
            # LLM returns structured JSON (SlidesSpec), not PptxGenJS code.
            # Validate via Pydantic, fix once if needed, then hand off to renderer.
            sanitized = _sanitize_json(raw)
            spec: SlidesSpec | None = None
            try:
                spec = SlidesSpec.model_validate_json(sanitized)
            except (ValidationError, Exception) as first_err:
                logging.warning(
                    "SlidesSpec validation failed for artifact %d — asking LLM to fix: %s",
                    artifact_id, str(first_err)[:300],
                )
                update_studio_artifact(artifact_id, progress_message="JSON 格式有誤，正在要求 AI 修正…")
                fixed_json = await _fix_slides_json(sanitized, str(first_err))
                try:
                    spec = SlidesSpec.model_validate_json(fixed_json)
                    sanitized = fixed_json
                except (ValidationError, Exception) as second_err:
                    update_studio_artifact(
                        artifact_id,
                        status="error",
                        content_json="{}",
                        content_text=raw,
                        error_message=f"簡報 JSON 驗證失敗，請稍後重試。({str(second_err)[:200]})",
                    )
                    return

            spec_json = spec.model_dump_json()
            update_studio_artifact(
                artifact_id,
                content_json=spec_json,
                content_text="",
                progress_message="AI 產出完成，正在建立簡報…",
            )
            _task = asyncio.create_task(_generate_slides_from_json(artifact_id, spec_json))

            def _log_task_exception(fut: asyncio.Future) -> None:
                if fut.exception():
                    logging.error(
                        "Background slides task failed for artifact=%d: %s",
                        artifact_id,
                        fut.exception(),
                    )

            _task.add_done_callback(_log_task_exception)
            return

        data = json.loads(raw)
        content_text = _format_text(artifact_type, data)

        content_json_str = json.dumps(data, ensure_ascii=False)
        update_studio_artifact(
            artifact_id,
            status="done",
            content_json=content_json_str,
            content_text=content_text,
        )

    except json.JSONDecodeError:
        # LLM output wasn't valid JSON — store raw text for debugging
        update_studio_artifact(
            artifact_id,
            status="error",
            content_json="{}",
            content_text=raw,
            error_message="AI 回應格式有誤，無法解析內容，請稍後重試。",
        )
    except Exception:
        logging.exception("Studio artifact generation failed: project=%s type=%s", project_id, artifact_type)
        update_studio_artifact(
            artifact_id,
            status="error",
            error_message="生成失敗，請確認文件內容有效後稍後重試。",
        )


_THUMB_ROOT = Path("/data/thumbnails")

# Limit concurrent slides renders (JSON renderer + LibreOffice thumbnail generation
# both consume significant CPU/memory; cap at 2 to avoid OOM on constrained hosts).
_SLIDES_SEMAPHORE = asyncio.Semaphore(2)


async def _generate_slides_from_json(artifact_id: int, spec_json: str) -> None:
    """
    Background task: convert a validated SlidesSpec JSON string into a PPTX
    file using the deterministic slides_renderer.js template engine, then
    generate thumbnails and run optional Vision QA.
    """
    async with _SLIDES_SEMAPHORE:
        await _generate_slides_from_json_inner(artifact_id, spec_json)


async def _generate_slides_from_json_inner(artifact_id: int, spec_json: str) -> None:
    import time
    from app.services.pptx_runner_service import execute_slides_json, RunResult
    from app.services.thumbnail_service import generate_thumbnails

    t0 = time.monotonic()

    with tempfile.TemporaryDirectory() as tmp:
        pptx_tmp = str(Path(tmp) / "slides.pptx")
        update_studio_artifact(artifact_id, progress_message="正在渲染簡報範本…")

        result, stderr = await execute_slides_json(spec_json, pptx_tmp)

        if result != RunResult.SUCCESS:
            logging.error(
                "slides_renderer failed for artifact %d (%s): %s",
                artifact_id, result, stderr,
            )
            update_studio_artifact(
                artifact_id,
                status="error",
                error_message="簡報渲染失敗，請稍後重試。",
            )
            return

        # Persist PPTX for download
        persistent_dir = _THUMB_ROOT / str(artifact_id)
        persistent_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = persistent_dir / "slides.pptx"
        shutil.copy2(pptx_tmp, pptx_path)

        # Optional: add AI-generated illustrations via ComfyUI (non-fatal if unavailable)
        from app.routers.settings import _runtime_settings as _rs_slides
        from app.services import comfyui_service
        if _rs_slides.comfyui_api_url and await comfyui_service.is_available():
            try:
                spec = SlidesSpec.model_validate_json(spec_json)
                await _add_illustrations_to_pptx(artifact_id, pptx_path, spec)
            except Exception:
                logging.exception("ComfyUI illustration step failed for artifact=%d — skipping", artifact_id)

        update_studio_artifact(artifact_id, progress_message="正在生成投影片縮圖…")
        try:
            _, slide_count = await asyncio.to_thread(generate_thumbnails, artifact_id, str(pptx_path))
        except Exception:
            logging.exception("Thumbnail generation failed: artifact=%d", artifact_id)
            update_studio_artifact(
                artifact_id,
                status="error",
                error_message="縮圖生成失敗，請稍後重試。",
            )
            return

        if slide_count > 0:
            update_studio_artifact(artifact_id, slide_count=slide_count)

        # Optional Vision QA — only runs when vision_model is configured.
        # When issues are found we attempt one automated re-render with an
        # LLM-assisted fix prompt before giving up and marking as done.
        from app.routers.settings import _runtime_settings
        from app.services.vision_qa import visual_qa_check, classify_issues, IssueType
        from app.services.thumbnail_service import get_thumbnail_urls

        if _runtime_settings.vision_model:
            thumb_urls = get_thumbnail_urls(artifact_id)
            thumb_paths = [_THUMB_ROOT / str(artifact_id) / Path(u).name for u in thumb_urls]
            update_studio_artifact(artifact_id, progress_message="正在進行視覺品質檢查…")
            issues = await visual_qa_check(
                thumb_paths,
                api_base_url=_runtime_settings.llm_api_base_url,
                api_key=_runtime_settings.llm_api_key,
                model=_runtime_settings.vision_model,
            )
            problem_count = sum(len(s.get("issues", [])) for s in issues)
            if problem_count > 0:
                logging.warning(
                    "Vision QA found %d issue(s) in artifact %d: %s",
                    problem_count, artifact_id, issues,
                )
                # ── Feedback loop: attempt one automated fix ──────────────
                issue_types = classify_issues(issues)
                fixed_spec_json = await _apply_qa_fixes(spec_json, issues, issue_types)
                if fixed_spec_json:
                    logging.info(
                        "Vision QA: retrying artifact=%d with auto-fixed spec (issues: %s)",
                        artifact_id, [t.value for t in issue_types],
                    )
                    update_studio_artifact(
                        artifact_id,
                        progress_message="偵測到視覺問題，正在自動修正並重新渲染…",
                    )
                    # Light retry: re-render PPTX, re-embed ComfyUI illustrations, re-generate thumbnails
                    try:
                        retry_pptx = str(Path(tmp) / "slides_retry.pptx")
                        result2, stderr2 = await execute_slides_json(fixed_spec_json, retry_pptx)
                        if result2 == RunResult.SUCCESS:
                            shutil.copy2(retry_pptx, pptx_path)
                            # Re-embed existing ComfyUI illustrations (already generated, just need re-embed)
                            if _rs_slides.comfyui_api_url:
                                try:
                                    from app.schemas.slides import SlidesSpec as _RetrySpec
                                    retry_spec = _RetrySpec.model_validate_json(fixed_spec_json)
                                    await asyncio.to_thread(
                                        _embed_images_in_pptx,
                                        pptx_path,
                                        [
                                            (idx, s)
                                            for idx, s in enumerate(retry_spec.slides)
                                            if getattr(s, "layout", "") in _ILLUSTRATABLE_LAYOUTS
                                        ][:_MAX_ILLUSTRATIONS],
                                        [
                                            (_COMFYUI_IMG_ROOT / str(artifact_id) / f"slide_{idx:03d}.png")
                                            if (_COMFYUI_IMG_ROOT / str(artifact_id) / f"slide_{idx:03d}.png").exists()
                                            else None
                                            for idx, s in enumerate(retry_spec.slides)
                                            if getattr(s, "layout", "") in _ILLUSTRATABLE_LAYOUTS
                                        ][:_MAX_ILLUSTRATIONS],
                                    )
                                    logging.info("Re-embedded ComfyUI illustrations after QA retry for artifact=%d", artifact_id)
                                except Exception:
                                    logging.exception("Failed to re-embed illustrations after QA retry for artifact=%d", artifact_id)
                            _, sc = await asyncio.to_thread(generate_thumbnails, artifact_id, str(pptx_path))
                            if sc > 0:
                                update_studio_artifact(artifact_id, slide_count=sc)
                            logging.info("Vision QA retry render succeeded for artifact=%d", artifact_id)
                        else:
                            logging.warning("Vision QA retry render failed for artifact=%d: %s", artifact_id, stderr2)
                    except Exception:
                        logging.exception("Vision QA retry failed for artifact=%d — accepting original", artifact_id)

        elapsed = time.monotonic() - t0
        logging.info(
            "slides_from_json artifact=%d done in %.1fs (renderer+thumbnails+qa)",
            artifact_id, elapsed,
        )
        update_studio_artifact(artifact_id, status="done", progress_message="")


# ── Vision QA auto-fix helpers ────────────────────────────────────────────────

_QA_FIX_PROMPT = """\
你是簡報文字修正助手。以下投影片有視覺品質問題，請修正**僅有問題的投影片**。

## 問題清單
{issue_summary}

## 修正規則
- low_contrast: 不修改文字（此問題由主題色碼處理）
- text_overflow: 將該頁所有文字欄位縮短 30%，保持語意完整
- excessive_blank: 在該頁增加 1-2 個 cards/items/points，用文件相關內容填充
- overlap: 減少該頁的 cards 或 steps 數量至 3 個以內

## 要求
- 只輸出修正後的完整 JSON（整份 slides spec）
- 保留未出問題投影片的原樣
- 輸出合法 JSON，不加 ```json 或說明文字
- 繁體中文

原始 JSON：
{spec_json}"""


async def _apply_qa_fixes(
    spec_json: str,
    issues: list[dict],
    issue_types: "set",
) -> str | None:
    """
    Ask LLM to fix only the problematic slides based on VLM feedback.
    Has a 60-second timeout to prevent infinite hangs.
    Returns fixed JSON string, or None if the fix attempt fails.
    """
    from app.services.vision_qa import IssueType

    # If only low_contrast issues, fix theme locally without LLM
    if issue_types == {IssueType.LOW_CONTRAST}:
        try:
            spec_dict = json.loads(spec_json)
            if spec_dict.get("theme") == "frost":
                spec_dict["theme"] = "ocean"
                logging.info("QA auto-fix: switched theme frost → ocean (low_contrast)")
            fixed = SlidesSpec.model_validate(spec_dict)
            return fixed.model_dump_json()
        except Exception:
            logging.exception("QA local theme fix failed")
            return None

    # Summarise issues for the prompt
    issue_lines: list[str] = []
    for slide_issue in issues:
        slide_num = slide_issue.get("slide", "?")
        for iss in slide_issue.get("issues", []):
            iss_type = iss.get("type", "unknown")
            iss_desc = iss.get("description", "")
            issue_lines.append(f"  - 第{slide_num}頁 [{iss_type}]: {iss_desc}")
    issue_summary = "\n".join(issue_lines) or "  (unspecified issues)"

    fix_prompt = _QA_FIX_PROMPT.format(issue_summary=issue_summary, spec_json=spec_json)
    messages = [ChatMessage(role=MessageRole.USER, content=fix_prompt)]
    fix_client = _fresh_async_client()
    llm = get_llm(async_client=fix_client)
    try:
        async def _stream_fix():
            parts: list[str] = []
            async for chunk in await llm.astream_chat(messages):
                if chunk.delta:
                    parts.append(chunk.delta)
            return "".join(parts).strip()

        # 60-second hard timeout
        raw = await asyncio.wait_for(_stream_fix(), timeout=60)
        fixed = _strip_code_fence(raw)

        # Validate the fixed JSON is still a valid SlidesSpec
        SlidesSpec.model_validate_json(fixed)
        return fixed
    except asyncio.TimeoutError:
        logging.warning("QA fix LLM call timed out after 60s — skipping retry")
        return None
    except Exception:
        logging.exception("QA fix LLM call failed for spec — skipping retry")
        return None
    finally:
        await fix_client.aclose()




# Maps SlidesSpec.theme → color palette description for ComfyUI Flux prompts
_THEME_COLOR_HINTS: dict[str, str] = {
    "tech":   "dark navy background with electric blue neon glow, cyan circuit patterns, futuristic digital aesthetic",
    "ocean":  "deep teal gradient with seafoam highlights, flowing water patterns, aquatic marine motifs",
    "golden": "rich dark brown background with warm amber and gold metallic accents, premium luxury feel",
    "frost":  "steel blue geometric shapes on soft gray background, ice crystal patterns, silver metallic accents, NOT white background",
    "garden": "warm ivory background with sage green botanical illustrations, earth tone organic shapes",
    "sports": "bold crimson red streaks on deep navy, dynamic motion lines, high-energy athletic aesthetic",
}

# Maps SlidesSpec.theme → illustration visual style
_THEME_STYLE_HINTS: dict[str, str] = {
    "tech":   "futuristic 3D render with glowing edges and depth",
    "ocean":  "soft watercolor illustration with flowing organic forms",
    "golden": "elegant art deco style with metallic textures and geometry",
    "frost":  "clean geometric illustration with crystalline angular details",
    "garden": "warm botanical watercolor with hand-drawn organic feel",
    "sports": "dynamic motion graphics with bold graphic design and energy",
}

# Layouts worth illustrating.
# Only layouts that have a natural empty zone for an image:
#   cover          – right-half panel (text is confined to left 60%)
#   section_divider – full-page background (image sent to back layer)
#   content_with_icon – icon+text on left, right third is empty
# Excluded: big_number / card_grid / dual_column / process_flow
#   (cards/columns fill the full slide width — any image would overlap content)
_ILLUSTRATABLE_LAYOUTS = frozenset(
    {"cover", "section_divider", "content_with_icon"}
)

# Max illustrations per presentation to avoid very long generation times
_MAX_ILLUSTRATIONS = 5

# ComfyUI images temp directory
_COMFYUI_IMG_ROOT = Path("/data/comfyui_images")


_SLIDE_ILLUSTRATION_BATCH_PROMPT = """\
You are an expert at creating image prompts for business presentation slides.
Generate English image prompts for Flux AI for each slide listed below.

CRITICAL RULES:
- All prompts must be 100% English — Flux does not support Chinese
- Each prompt must visually represent its SPECIFIC topic using concrete objects — NOT generic nature or plain backgrounds
- Topic-to-visual mapping (use as inspiration):
    AI / machine learning  → neural network nodes, circuit board traces, data stream particles
    compliance / legal     → official seal, balance scale, document stack, gavel, protective shield
    risk / security        → layered firewall, lock mechanism, warning indicator, armored vault
    data / analytics       → 3D bar chart, holographic dashboard, flowing data nodes, graph lines
    process / workflow     → connected gear system, flowchart pipeline, interlocking steps
    finance / investment   → ascending chart bars, coin stack, growth curve, currency symbols
    military / defense     → insignia badge, strategic map, command center, structured ranks
    education / research   → open book, magnifying glass, academic structure, knowledge graph
    health / medicine      → cross symbol, cellular structure, diagnostic display, clean lab
    environment / energy   → leaf circuit, solar array, wind turbine, green gradient sphere
- Interpret ANY topic not listed above by identifying its core visual metaphors

Visual style for ALL prompts: {style_hint}
Color palette for ALL prompts: {color_hint}
Constraints: No human faces; no text or letters in the image; landscape 4:3 format

SLIDES:
{slides_list}

Return a JSON array of strings — one prompt per slide (30-60 words each), in the same order:
["prompt for slide 1", "prompt for slide 2", ...]
Return ONLY the JSON array, nothing else."""


async def _generate_image_prompts_for_slides(
    eligible: list[tuple[int, object]],
    theme: str = "tech",
) -> list[str | None]:
    """
    Use ONE LLM call to generate English image prompts for all eligible slides.
    Returns a list of prompt strings (or None if generation fails per slide).
    """
    if not eligible:
        return []

    color_hint = _THEME_COLOR_HINTS.get(theme, _THEME_COLOR_HINTS["tech"])
    style_hint = _THEME_STYLE_HINTS.get(theme, _THEME_STYLE_HINTS["tech"])

    # Build the slides list section
    slide_lines: list[str] = []
    for i, (_slide_idx, slide_data) in enumerate(eligible):
        title = getattr(slide_data, "title", "") or getattr(slide_data, "quote", "") or ""
        layout = getattr(slide_data, "layout", "")
        hint_parts: list[str] = []
        subtitle = getattr(slide_data, "subtitle", "") or getattr(slide_data, "description", "")
        if subtitle:
            hint_parts.append(subtitle)
        cards: list = getattr(slide_data, "cards", []) or []
        if cards:
            hint_parts.extend(getattr(c, "title", "") for c in cards[:3])
        content_hint = "; ".join(p for p in hint_parts if p) or "(none)"
        slide_lines.append(f"Slide {i + 1}: layout={layout}, title={title}, content={content_hint}")

    slides_list = "\n".join(slide_lines)
    user_msg = _SLIDE_ILLUSTRATION_BATCH_PROMPT.format(
        style_hint=style_hint,
        color_hint=color_hint,
        slides_list=slides_list,
    )

    client = _fresh_async_client()
    llm = get_llm(async_client=client)
    try:
        parts: list[str] = []
        async for chunk in await llm.astream_chat(
            [ChatMessage(role=MessageRole.USER, content=user_msg)]
        ):
            if chunk.delta:
                parts.append(chunk.delta)
        raw = "".join(parts).strip()
        raw = _strip_code_fence(raw)

        parsed = json.loads(raw)
        if isinstance(parsed, list):
            # Pad or truncate to match eligible count
            result: list[str | None] = []
            for i in range(len(eligible)):
                val = parsed[i] if i < len(parsed) else None
                result.append(str(val).strip() if val else None)
            return result
        raise ValueError(f"Expected JSON array, got: {type(parsed)}")

    except Exception:
        logging.exception("Batch image prompt generation failed — no illustrations")
        return [None] * len(eligible)
    finally:
        await client.aclose()



def _embed_images_in_pptx(
    pptx_path: Path,
    slides_info: list[tuple[int, object]],
    image_paths: list[Path | None],
) -> None:
    """
    Open the PPTX with python-pptx and embed generated illustrations.
    Positions vary by layout:
      - cover: right half large image
      - section_divider: full-page background (sent to back)
      - others: right-side inset
    Non-fatal: skips individual images that fail.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(str(pptx_path))
    modified = False

    for (slide_idx, slide_data), img_path in zip(slides_info, image_paths):
        if img_path is None or not img_path.exists():
            continue
        if slide_idx >= len(prs.slides):
            continue

        sld = prs.slides[slide_idx]
        layout = getattr(slide_data, "layout", "")

        try:
            if layout == "section_divider":
                # Full-page background — add picture then send to back
                pic = sld.shapes.add_picture(str(img_path), Inches(0), Inches(0), Inches(10), Inches(5.625))
                sp = pic._element
                sld.shapes._spTree.remove(sp)
                sld.shapes._spTree.insert(2, sp)  # behind all other elements
            elif layout == "cover":
                # Right-panel illustration — starts at 60% width, full height
                # Title text is left-aligned so stays within the left 60%
                sld.shapes.add_picture(str(img_path), Inches(6.0), Inches(0), Inches(4.0), Inches(5.625))
            else:
                # content_with_icon: icon+text block is on the left, right third is empty
                sld.shapes.add_picture(str(img_path), Inches(7.2), Inches(1.5), Inches(2.4), Inches(2.4))
            modified = True
            logging.info("Embedded illustration for slide %d (layout=%s)", slide_idx, layout)
        except Exception:
            logging.exception("Failed to embed image for slide %d", slide_idx)

    if modified:
        prs.save(str(pptx_path))
        logging.info("PPTX saved with illustrations: %s", pptx_path)


async def _add_illustrations_to_pptx(
    artifact_id: int,
    pptx_path: Path,
    spec: object,
) -> None:
    """
    Generate ComfyUI illustrations for eligible slides and embed them into the PPTX.
    Non-fatal: any failure is logged and silently skipped.
    """
    from app.services.comfyui_service import generate_image as comfyui_generate

    slides = getattr(spec, "slides", [])
    eligible = [
        (i, s) for i, s in enumerate(slides)
        if getattr(s, "layout", "") in _ILLUSTRATABLE_LAYOUTS
    ][:_MAX_ILLUSTRATIONS]

    if not eligible:
        logging.info("No illustratable slides found for artifact=%d", artifact_id)
        return

    logging.info("Generating ComfyUI illustrations for %d slides (artifact=%d)", len(eligible), artifact_id)
    update_studio_artifact(artifact_id, progress_message=f"正在生成投影片插圖 (0/{len(eligible)})…")

    # Step 1: Generate English prompts via LLM (pass theme for color palette guidance)
    theme = getattr(spec, "theme", "tech") or "tech"
    img_prompts = await _generate_image_prompts_for_slides(eligible, theme=theme)

    # Step 2: Generate images via ComfyUI (max 2 concurrent)
    img_dir = _COMFYUI_IMG_ROOT / str(artifact_id)
    img_dir.mkdir(parents=True, exist_ok=True)

    sem = asyncio.Semaphore(2)
    image_paths: list[Path | None] = [None] * len(eligible)
    completed = [0]

    async def gen_one(idx: int, prompt: str | None) -> None:
        if not prompt:
            return
        async with sem:
            save_path = img_dir / f"slide_{eligible[idx][0]:03d}.png"
            ok = await comfyui_generate(prompt, save_path)
            if ok:
                image_paths[idx] = save_path
            completed[0] += 1
            update_studio_artifact(
                artifact_id,
                progress_message=f"正在生成投影片插圖 ({completed[0]}/{len(eligible)})…",
            )

    await asyncio.gather(*[gen_one(i, p) for i, p in enumerate(img_prompts)])

    success_count = sum(1 for p in image_paths if p is not None)
    if success_count == 0:
        logging.warning("ComfyUI: no illustrations generated for artifact=%d", artifact_id)
        return

    # Step 3: Embed images into PPTX
    await asyncio.to_thread(_embed_images_in_pptx, pptx_path, eligible, image_paths)
    logging.info(
        "ComfyUI illustrations: %d/%d embedded for artifact=%d",
        success_count, len(eligible), artifact_id,
    )
