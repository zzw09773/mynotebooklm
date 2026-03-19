"""
Server-side PPTX generation using python-pptx.

Implements the same 10 themes and 12 layout types as the frontend PptxGenJS
output (SlidesViewer.tsx).  Used for the soffice → fitz thumbnail pipeline.

Canvas: LAYOUT_16x9 — 10" × 5.625"
Colors: 6-char hex, no '#' prefix (consistent with PptxGenJS convention)
"""
import json
import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

log = logging.getLogger(__name__)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ── Theme palette (mirrors SLIDE_THEMES in SlidesViewer.tsx) ──

THEMES: dict[str, dict[str, str]] = {
    "tech-innovation":   dict(bg="1E1E1E", accent="0066FF", title="FFFFFF", text="CCCCCC", muted="888888", card="2A2A2A"),
    "midnight-galaxy":   dict(bg="2B1E3E", accent="A490C2", title="E6E6FA", text="C8B8E0", muted="6B5B8A", card="3A2D50"),
    "ocean-depths":      dict(bg="1A2332", accent="2D8B8B", title="F1FAEE", text="A8DADC", muted="5A8A8A", card="243040"),
    "modern-minimalist": dict(bg="FFFFFF", accent="708090", title="36454F", text="36454F", muted="A0A0A0", card="F0F0F0"),
    "sunset-boulevard":  dict(bg="264653", accent="E76F51", title="E9C46A", text="F4A261", muted="A09060", card="314D5E"),
    "forest-canopy":     dict(bg="2D4A2B", accent="A4AC86", title="FAF9F6", text="C8CCB8", muted="7D8471", card="3A5C38"),
    "golden-hour":       dict(bg="4A403A", accent="F4A900", title="D4B896", text="D4B896", muted="C1666B", card="5A4E47"),
    "arctic-frost":      dict(bg="FAFAFA", accent="4A6FA5", title="1A2332", text="334155", muted="909090", card="D4E4F7"),
    "desert-rose":       dict(bg="E8D5C4", accent="B87D6D", title="5D2E46", text="5D2E46", muted="D4A5A5", card="F0E4D8"),
    "botanical-garden":  dict(bg="F5F3ED", accent="4A7C59", title="333333", text="555555", muted="B7472A", card="EBE9E1"),
}

_HEX_RE = set("0123456789ABCDEFabcdef")


# ── Low-level helpers ──────────────────────────────────────────

def _rgb(hex6: str) -> RGBColor:
    h = hex6.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_bg(slide, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color_hex)


def _add_rect(slide, x_in: float, y_in: float, w_in: float, h_in: float, color_hex: str) -> None:
    """Add a solid filled rectangle (autoshape type 1 = rectangle)."""
    shape = slide.shapes.add_shape(
        1,  # MSO autoshape: rectangle
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color_hex)
    shape.line.color.rgb = _rgb(color_hex)  # same as fill → invisible border


def _add_text(
    slide,
    text: str,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    *,
    size_pt: int,
    color_hex: str,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.color.rgb = _rgb(color_hex)
    run.font.bold = bold
    run.font.italic = italic


def _add_bullets(
    slide,
    bullets: list[str],
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    *,
    size_pt: int,
    color_hex: str,
) -> None:
    if not bullets:
        return
    txBox = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.size = Pt(size_pt)
        p.font.color.rgb = _rgb(color_hex)
        p.space_after = Pt(6)


# ── Layout builders ───────────────────────────────────────────

def _build_cover(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0.35, 1.3, 0.07, 3.0, t["accent"])
    for size, px, py in [[0.38, 9.05, 4.45], [0.24, 9.28, 4.68], [0.12, 9.45, 4.85]]:
        _add_rect(slide, px, py, size, size, t["accent"])
    _add_text(slide, s.get("title", ""), 0.7, 1.6, 7.9, 1.4, size_pt=36, color_hex=t["title"], bold=True)
    if s.get("subtitle"):
        _add_text(slide, s["subtitle"], 0.7, 3.2, 7.9, 0.65, size_pt=18, color_hex=t["muted"])


def _build_content(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0, 0, 10, 0.09, t["accent"])
    _add_text(slide, s.get("title", ""), 0.5, 0.15, 9, 0.65, size_pt=22, color_hex=t["title"], bold=True)
    bullet_y = 1.08 if s.get("subtitle") else 0.92
    if s.get("subtitle"):
        _add_text(slide, s["subtitle"], 0.5, 0.85, 9, 0.35, size_pt=12, color_hex=t["muted"])
    _add_bullets(
        slide, s.get("bullets", []),
        0.5, bullet_y, 9, 5.625 - bullet_y - 0.45,
        size_pt=13, color_hex=t["text"],
    )


def _build_section(slide, s: dict, t: dict, index: int) -> None:
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0, 0, 10, 2.55, t["accent"])
    _add_text(slide, f"SECTION {index + 1:02d}", 0.6, 0.45, 8.5, 0.45, size_pt=11, color_hex="FFFFFF")
    _add_text(slide, s.get("title", ""), 0.6, 1.0, 8.5, 1.3, size_pt=30, color_hex="FFFFFF", bold=True)
    if s.get("subtitle"):
        _add_text(slide, s["subtitle"], 0.6, 2.9, 8.5, 0.85, size_pt=16, color_hex=t["text"])


def _build_conclusion(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0.35, 1.0, 0.07, 1.5, t["accent"])
    _add_text(slide, s.get("title", ""), 0.7, 1.05, 8.5, 1.1, size_pt=30, color_hex=t["title"], bold=True)
    if s.get("subtitle"):
        _add_text(slide, s["subtitle"], 0.7, 2.25, 8.5, 0.55, size_pt=15, color_hex=t["muted"])
    # Key-point chips (simplified: just bullets)
    chips = s.get("bullets", [])[:4]
    if chips:
        chip_y = 3.15 if s.get("subtitle") else 2.9
        chip_w, chip_h = 4.2, 0.55
        for ci, chip_text in enumerate(chips):
            cx = 0.5 + (ci % 2) * (chip_w + 0.2)
            cy = chip_y + (ci // 2) * (chip_h + 0.15)
            _add_rect(slide, cx, cy, chip_w, chip_h, t["card"])
            _add_rect(slide, cx, cy, 0.05, chip_h, t["accent"])
            colon_idx = next((j for j, c in enumerate(chip_text) if c in "：:"), -1)
            label = chip_text[:colon_idx] if 0 < colon_idx <= 20 else chip_text[:25]
            _add_text(slide, label, cx + 0.15, cy + 0.08, chip_w - 0.25, chip_h - 0.1, size_pt=12, color_hex=t["text"])


def _build_quote(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["bg"])
    _add_text(slide, "\u201C", 0.3, 0.1, 1.5, 1.5, size_pt=80, color_hex=t["accent"], bold=True)
    bullets = s.get("bullets", [])
    quote = bullets[0] if bullets else s.get("title", "")
    _add_text(slide, quote, 1.0, 1.1, 7.9, 2.6, size_pt=22, color_hex=t["title"], italic=True, align=PP_ALIGN.CENTER)
    if s.get("subtitle"):
        _add_text(slide, f"\u2014 {s['subtitle']}", 1.0, 4.3, 8.7, 0.5, size_pt=13, color_hex=t["muted"], align=PP_ALIGN.RIGHT)


def _build_stats(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0, 0, 10, 0.09, t["accent"])
    _add_text(slide, s.get("title", ""), 0.5, 0.15, 9, 0.65, size_pt=22, color_hex=t["title"], bold=True)
    cols = [0.35, 5.2]
    rows = [1.05, 3.1]
    for ci, bullet_text in enumerate(s.get("bullets", [])[:4]):
        cx, cy = cols[ci % 2], rows[ci // 2]
        _add_rect(slide, cx, cy, 4.5, 1.7, t["card"])
        _add_rect(slide, cx, cy, 4.5, 0.05, t["accent"])
        colon_idx = next((j for j, c in enumerate(bullet_text) if c in "：:"), -1)
        if 0 < colon_idx <= 8:
            big_text = bullet_text[:colon_idx]
            small_text = bullet_text[colon_idx + 1:].strip()
        else:
            big_text = bullet_text[:5]
            small_text = bullet_text[5:]
        _add_text(slide, big_text, cx + 0.2, cy + 0.12, 4.1, 0.65, size_pt=24, color_hex=t["accent"], bold=True)
        if small_text:
            _add_text(slide, small_text, cx + 0.2, cy + 0.82, 4.1, 0.75, size_pt=12, color_hex=t["text"])


# ── New layout builders ───────────────────────────────────────

def _build_big_number(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0, 0, 10, 0.09, t["accent"])
    _add_text(slide, s.get("title", ""), 0.5, 0.2, 9, 0.6, size_pt=20, color_hex=t["title"], bold=True)
    metric = s.get("metric", "—")
    unit = s.get("unit", "")
    metric_str = f"{metric} {unit}".strip()
    _add_text(slide, metric_str, 1.5, 1.1, 7.0, 2.4, size_pt=72, color_hex=t["accent"], bold=True, align=PP_ALIGN.CENTER)
    _add_rect(slide, 3.5, 3.6, 3.0, 0.04, t["muted"])
    if s.get("label"):
        _add_text(slide, s["label"], 1.5, 3.75, 7.0, 0.65, size_pt=18, color_hex=t["text"], align=PP_ALIGN.CENTER)
    bullets = s.get("bullets", [])
    if bullets:
        _add_text(slide, bullets[0], 1.5, 4.55, 7.0, 0.65, size_pt=12, color_hex=t["muted"], align=PP_ALIGN.CENTER)


def _build_dual_card(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0, 0, 10, 0.09, t["accent"])
    _add_text(slide, s.get("title", ""), 0.5, 0.18, 9, 0.6, size_pt=20, color_hex=t["title"], bold=True)
    for ci, card_key in enumerate(["left_card", "right_card"]):
        card = s.get(card_key)
        if not card:
            continue
        cx = 0.3 if ci == 0 else 5.2
        _add_rect(slide, cx, 0.95, 4.5, 4.3, t["card"])
        _add_rect(slide, cx, 0.95, 4.5, 0.06, t["accent"])
        _add_text(slide, card.get("title", ""), cx + 0.2, 1.1, 4.1, 0.55, size_pt=15, color_hex=t["title"], bold=True)
        _add_bullets(slide, card.get("bullets", []), cx + 0.2, 1.75, 4.1, 3.3, size_pt=12, color_hex=t["text"])


def _build_multi_card(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0, 0, 10, 0.09, t["accent"])
    _add_text(slide, s.get("title", ""), 0.5, 0.18, 9, 0.6, size_pt=20, color_hex=t["title"], bold=True)
    cards = (s.get("cards") or [])[:6]
    cols = 3
    card_w = 2.8
    card_h = 1.95 if len(cards) > 3 else 3.5
    start_x = 0.25
    start_y = 0.95
    for ci, card in enumerate(cards):
        col = ci % cols
        row = ci // cols
        cx = start_x + col * (card_w + 0.25)
        cy = start_y + row * (card_h + 0.2)
        _add_rect(slide, cx, cy, card_w, card_h, t["card"])
        _add_rect(slide, cx, cy, card_w, 0.05, t["accent"])
        _add_text(slide, "◆", cx + 0.15, cy + 0.15, 0.5, 0.4, size_pt=14, color_hex=t["accent"])
        _add_text(slide, card.get("title", ""), cx + 0.15, cy + 0.55, card_w - 0.3, 0.5, size_pt=13, color_hex=t["title"], bold=True)
        if card.get("description"):
            _add_text(slide, card["description"], cx + 0.15, cy + 1.1, card_w - 0.3, card_h - 1.25, size_pt=11, color_hex=t["text"])


def _build_table(slide, s: dict, t: dict) -> None:
    from pptx.util import Pt as _Pt
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0, 0, 10, 0.09, t["accent"])
    _add_text(slide, s.get("title", ""), 0.5, 0.18, 9, 0.6, size_pt=20, color_hex=t["title"], bold=True)
    headers = s.get("headers") or []
    rows_data = (s.get("rows") or [])[:8]
    if not headers:
        return
    col_count = len(headers)
    row_count = 1 + len(rows_data)
    tbl = slide.shapes.add_table(row_count, col_count, Inches(0.4), Inches(1.0), Inches(9.2), Inches(4.3)).table
    col_w_emu = int(Inches(9.2) / col_count)
    for ci in range(col_count):
        tbl.columns[ci].width = col_w_emu
    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(t["accent"])
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = _Pt(12)
        para.font.color.rgb = _rgb("FFFFFF")
        para.alignment = PP_ALIGN.CENTER
    # Data rows
    for ri, row_vals in enumerate(rows_data):
        bg = t["card"] if ri % 2 == 0 else t["bg"]
        for ci in range(col_count):
            cell = tbl.cell(ri + 1, ci)
            cell.text = (row_vals[ci] if ci < len(row_vals) else "")
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(bg)
            para = cell.text_frame.paragraphs[0]
            para.font.size = _Pt(11)
            para.font.color.rgb = _rgb(t["text"])


def _build_flow(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["bg"])
    _add_rect(slide, 0, 0, 10, 0.09, t["accent"])
    _add_text(slide, s.get("title", ""), 0.5, 0.18, 9, 0.6, size_pt=20, color_hex=t["title"], bold=True)
    steps = (s.get("steps") or [])[:5]
    if not steps:
        return
    step_w = 9.0 / len(steps)
    circle_r = 0.38
    circle_y = 1.6
    for si, step in enumerate(steps):
        cx = 0.5 + si * step_w + step_w / 2 - circle_r
        # Circle (oval autoshape = 9)
        oval = slide.shapes.add_shape(9, Inches(cx), Inches(circle_y), Inches(circle_r * 2), Inches(circle_r * 2))
        oval.fill.solid()
        oval.fill.fore_color.rgb = _rgb(t["accent"])
        oval.line.color.rgb = _rgb(t["accent"])
        _add_text(slide, str(si + 1), cx, circle_y, circle_r * 2, circle_r * 2, size_pt=14, color_hex="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        # Arrow connector (except last)
        if si < len(steps) - 1:
            arrow_x = cx + circle_r * 2 + 0.05
            arrow_w = step_w - circle_r * 2 - 0.1
            _add_rect(slide, arrow_x, circle_y + circle_r - 0.02, arrow_w, 0.04, t["muted"])
        # Label
        label_x = cx - 0.35
        label_w = circle_r * 2 + 0.7
        _add_text(slide, step.get("label", ""), label_x, circle_y + circle_r * 2 + 0.15, label_w, 0.5, size_pt=12, color_hex=t["title"], bold=True, align=PP_ALIGN.CENTER)
        if step.get("description"):
            _add_text(slide, step["description"], label_x, circle_y + circle_r * 2 + 0.75, label_w, 1.8, size_pt=10, color_hex=t["text"], align=PP_ALIGN.CENTER)


def _build_hero_text(slide, s: dict, t: dict) -> None:
    _set_bg(slide, t["accent"])
    # Corner decorations
    _add_rect(slide, 0, 0, 1.5, 0.08, "FFFFFF")
    _add_rect(slide, 0, 0, 0.08, 1.5, "FFFFFF")
    _add_rect(slide, 8.5, 5.545, 1.5, 0.08, "FFFFFF")
    _add_rect(slide, 9.92, 4.1, 0.08, 1.5, "FFFFFF")
    _add_text(slide, s.get("title", ""), 0.7, 1.3, 8.6, 1.8, size_pt=38, color_hex="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    if s.get("subtitle"):
        _add_text(slide, s["subtitle"], 0.7, 3.3, 8.6, 1.0, size_pt=18, color_hex="FFFFFF", align=PP_ALIGN.CENTER)


# ── Public entry point ────────────────────────────────────────

def generate_pptx(content_json: str, output_path: Path) -> None:
    """Generate a themed PPTX file from slides JSON and write to output_path."""
    data = json.loads(content_json)

    theme_name = data.get("theme", "tech-innovation")
    theme = dict(THEMES.get(theme_name, THEMES["tech-innovation"]))

    accent_override = data.get("accent_color", "")
    if accent_override and len(accent_override) == 6 and all(c in _HEX_RE for c in accent_override):
        theme["accent"] = accent_override

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank layout — no placeholders

    slides = data.get("slides", [])
    for i, s in enumerate(slides):
        sld = prs.slides.add_slide(blank_layout)
        layout = s.get("layout_type", "content")

        if layout == "cover":
            _build_cover(sld, s, theme)
        elif layout == "stats":
            _build_stats(sld, s, theme)
        elif layout == "quote":
            _build_quote(sld, s, theme)
        elif layout == "section":
            _build_section(sld, s, theme, i)
        elif layout == "conclusion":
            _build_conclusion(sld, s, theme)
        elif layout == "big_number":
            _build_big_number(sld, s, theme)
        elif layout == "dual_card":
            _build_dual_card(sld, s, theme)
        elif layout == "multi_card":
            _build_multi_card(sld, s, theme)
        elif layout == "table":
            _build_table(sld, s, theme)
        elif layout == "flow":
            _build_flow(sld, s, theme)
        elif layout == "hero_text":
            _build_hero_text(sld, s, theme)
        else:
            _build_content(sld, s, theme)

    prs.save(str(output_path))
    log.info("PPTX written: %s (%d slides)", output_path, len(slides))
